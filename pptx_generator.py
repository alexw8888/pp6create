#!/usr/bin/env python3
"""
PowerPoint Generator Module
Creates PowerPoint presentations from source materials
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import os
import math
from pathlib import Path
from typing import List, Dict, Tuple, Optional
import json
from PIL import Image
from dotenv import load_dotenv

# Load environment variables
load_dotenv()


class PPTXGenerator:
    def __init__(self, width=1024, height=768):
        self.width = width
        self.height = height
        
        # Default settings (can be overridden)
        self.font_size = int(os.getenv('PPTX_FONT_SIZE', os.getenv('FONT_SIZE', '40')))
        self.font_family = os.getenv('FONT_FAMILY', 'Arial')
        self.font_color = self._parse_color(os.getenv('FONT_COLOR', '0xFFFFFF'))
        self.top_margin = int(os.getenv('TOP_MARGIN', '100'))
        self.page_break_every = int(os.getenv('PAGE_BREAK_EVERY', '4'))
        
        # Shadow settings
        self.add_text_shadow = os.getenv('ADD_TEXT_SHADOW', 'true').lower() == 'true'
        self.shadow_offset_x = int(os.getenv('SHADOW_OFFSET_X', '2'))
        self.shadow_offset_y = int(os.getenv('SHADOW_OFFSET_Y', '2'))
        self.shadow_blur_radius = float(os.getenv('SHADOW_BLUR_RADIUS', '3'))
        
        # Create presentation
        self.prs = Presentation()
        
        # Set slide size (convert pixels to inches - PowerPoint uses 96 DPI)
        self.prs.slide_width = Inches(self.width / 96.0)
        self.prs.slide_height = Inches(self.height / 96.0)
    
    def _parse_color(self, color_str: str) -> Tuple[int, int, int]:
        """Parse color from hex string to RGB tuple."""
        if color_str.startswith('0x'):
            color_str = color_str[2:]
        elif color_str.startswith('#'):
            color_str = color_str[1:]
        
        # Convert hex to RGB
        color_int = int(color_str, 16)
        r = (color_int >> 16) & 255
        g = (color_int >> 8) & 255
        b = color_int & 255
        return (r, g, b)
    
    def _add_text_shadow_to_run(self, run):
        """Add text shadow to a text run using XML manipulation."""
        if not self.add_text_shadow:
            return
        
        try:
            # Access the run's XML element
            run_element = run._r
            
            # Get or create the rPr (run properties) element
            rPr = run_element.get_or_add_rPr()
            
            # Check if shadow already exists
            existing_effectLst = rPr.find(qn('a:effectLst'))
            if existing_effectLst is not None:
                return
            
            # Create the effectLst element
            effectLst = etree.SubElement(rPr, qn('a:effectLst'))
            
            # Calculate shadow parameters
            blur_radius = int(self.shadow_blur_radius * 12700)  # Convert points to EMUs
            dist_x = self.shadow_offset_x * 12700
            dist_y = self.shadow_offset_y * 12700
            
            # Calculate distance and angle from x,y offsets
            distance = int(math.sqrt(dist_x**2 + dist_y**2))
            if dist_x == 0:
                angle = 5400000 if dist_y > 0 else 16200000  # 90 or 270 degrees
            else:
                angle_rad = math.atan2(dist_y, dist_x)
                angle_deg = math.degrees(angle_rad)
                angle = int(angle_deg * 60000) % 21600000
            
            # Create the outerShdw element
            outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
            outerShdw.set('blurRad', str(blur_radius))
            outerShdw.set('dist', str(distance))
            outerShdw.set('dir', str(angle))
            
            # Add shadow color (black)
            srgbClr = etree.SubElement(outerShdw, qn('a:srgbClr'))
            srgbClr.set('val', '000000')
            
            # Add transparency (60% opacity for shadow)
            alpha = etree.SubElement(srgbClr, qn('a:alpha'))
            alpha.set('val', '60000')
            
        except Exception as e:
            print(f"Warning: Could not add text shadow: {e}")
    
    def add_slide_with_background(self, background_path: str = None, text: str = None, 
                                 position: Dict = None, font_config: Dict = None,
                                 text_align: str = 'center', background_color: str = '#000000'):
        """Add a slide with background image/color and optional text."""
        # Use blank slide layout
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add background color first
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, self.prs.slide_height
        )
        bg_shape.fill.solid()
        
        # Parse background color
        bg_color_rgb = self._parse_color(background_color)
        bg_shape.fill.fore_color.rgb = RGBColor(*bg_color_rgb)
        bg_shape.line.color.rgb = RGBColor(*bg_color_rgb)
        
        # Load and add image with preserved aspect ratio
        if background_path and os.path.exists(background_path):
            with Image.open(background_path) as img:
                img_width, img_height = img.size
                
            # Calculate aspect ratios
            slide_aspect = self.prs.slide_width / self.prs.slide_height
            img_aspect = img_width / img_height
            
            # Calculate dimensions to preserve aspect ratio
            if img_aspect > slide_aspect:
                # Image is wider - fit to width
                new_width = self.prs.slide_width
                new_height = self.prs.slide_width / img_aspect
                left = Inches(0)
                top = (self.prs.slide_height - new_height) / 2
            else:
                # Image is taller - fit to height
                new_height = self.prs.slide_height
                new_width = self.prs.slide_height * img_aspect
                left = (self.prs.slide_width - new_width) / 2
                top = Inches(0)
            
            # Add the image
            pic = slide.shapes.add_picture(background_path, left, top, 
                                         width=new_width, height=new_height)
            
            # Send picture to back (but in front of black background)
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(3, pic._element)
        
        # Add text if provided
        if text:
            # Use position from JSON or default
            if position:
                # Convert pixel coordinates to inches
                text_left = Inches(position.get('x', 0) / 96.0)
                text_top = Inches(position.get('y', 100) / 96.0)
                text_width = Inches(position.get('width', 400) / 96.0)
                text_height = Inches(position.get('height', 100) / 96.0)
            else:
                # Default centered position
                text_width = self.prs.slide_width - Inches(1)
                text_left = Inches(0.5)
                text_top = Inches(self.top_margin / 72.0)
                text_height = self.prs.slide_height - text_top
            
            # Add text box
            text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
            text_frame = text_box.text_frame
            text_frame.clear()
            
            # Set text frame properties
            text_frame.margin_bottom = Inches(0.08)
            text_frame.margin_left = Inches(0)
            text_frame.margin_right = Inches(0)
            text_frame.margin_top = Inches(0)
            text_frame.word_wrap = True
            
            # Add text
            p = text_frame.paragraphs[0]
            
            # Set alignment based on parameter
            if text_align == 'left':
                p.alignment = PP_ALIGN.LEFT
            elif text_align == 'right':
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.CENTER
            
            run = p.add_run()
            run.text = text
            
            # Apply font configuration
            font = run.font
            if font_config:
                font.name = font_config.get('font_name', self.font_family)
                font.size = Pt(font_config.get('font_size', self.font_size))
                if 'font_color' in font_config:
                    font.color.rgb = RGBColor(*self._parse_color(font_config['font_color']))
                else:
                    font.color.rgb = RGBColor(*self.font_color)
            else:
                font.name = self.font_family
                font.size = Pt(self.font_size)
                font.color.rgb = RGBColor(*self.font_color)
            
            # Add text shadow if enabled
            self._add_text_shadow_to_run(run)
        
        return slide
    
    def generate_from_json_directory(self, source_dir: str) -> Presentation:
        """Generate PowerPoint using unified JSON/media processing rules."""
        source_path = Path(source_dir)
        
        # Collect all files
        json_files = sorted(source_path.glob('*.json'), key=lambda x: x.name)
        image_files = []
        video_files = []
        
        # Collect image files
        for ext in ['*.png', '*.jpg', '*.jpeg']:
            image_files.extend(source_path.glob(ext))
        
        # Collect video files
        video_files.extend(source_path.glob('*.mp4'))
        
        # Combine all media files
        all_media_files = sorted(image_files + video_files, key=lambda x: x.name)
        
        # Create a set of all unique base names
        all_base_names = set()
        for f in json_files:
            all_base_names.add(f.stem)
        for f in all_media_files:
            all_base_names.add(f.stem)
        
        # Process each unique base name
        for base_name in sorted(all_base_names):
            # Check for JSON file
            json_file = source_path / f"{base_name}.json"
            has_json = json_file.exists()
            
            # Check for media file
            media_file = None
            for ext in ['.png', '.jpg', '.jpeg', '.mp4']:
                potential_media = source_path / f"{base_name}{ext}"
                if potential_media.exists():
                    media_file = str(potential_media)
                    break
            
            if has_json:
                # Load JSON configuration
                with open(json_file, 'r', encoding='utf-8') as f:
                    config = json.load(f)
                
                # Extract configuration
                text = config.get('text', '')
                
                # Calculate PowerPoint-specific position with offsets
                position = None
                if 'x' in config and 'y' in config:
                    x = config.get('x', 0)
                    y = config.get('y', 0)
                    
                    # Apply PowerPoint-specific offsets if provided
                    x += config.get('pptxXoffset', 0)
                    y += config.get('pptxYoffset', 0)
                    
                    position = {
                        'x': x,
                        'y': y,
                        'width': config.get('width', 400),
                        'height': config.get('height', 100)
                    }
                
                # Calculate font size using scale factor
                base_font_size = config.get('fontSize', self.font_size)
                pptx_font_scale = config.get('pptxFontScale', 1.0)
                calculated_font_size = int(base_font_size * pptx_font_scale)
                
                font_config = {
                    'font_size': calculated_font_size,
                    'font_name': config.get('fontFamily', self.font_family),
                    'font_color': config.get('color', f'0x{self.font_color[0]:02X}{self.font_color[1]:02X}{self.font_color[2]:02X}')
                }
                
                # Get background color
                background_color = config.get('backgroundColor', '#000000')
                
                # Add slide with left alignment for JSON-based slides
                self.add_slide_with_background(media_file, text, position, font_config, 
                                             text_align='left', background_color=background_color)
            else:
                # No JSON - image only slide
                if media_file:
                    self.add_slide_with_background(media_file)
        
        return self.prs
    
    def generate_from_directory(self, source_dir: str) -> Presentation:
        """Generate PowerPoint from directory using unified rules."""
        source_path = Path(source_dir)
        
        # Check if directory has JSON files
        json_files = list(source_path.glob('*.json'))
        
        if json_files:
            # Use unified JSON/media processing
            return self.generate_from_json_directory(source_dir)
        
        # Legacy processing - check for song files
        text_files = list(source_path.glob('*.txt'))
        image_files = []
        
        # Collect image files
        for ext in ['*.png', '*.jpg', '*.jpeg']:
            image_files.extend(source_path.glob(ext))
        
        # Sort files
        text_files.sort(key=lambda x: x.name)
        image_files.sort(key=lambda x: x.name)
        
        # Check if any text file is a song
        song_file = None
        for txt_file in text_files:
            with open(txt_file, 'r', encoding='utf-8') as f:
                content = f.read()
                if 'Arrangement' in content:
                    song_file = txt_file
                    break
        
        if song_file:
            # Process as song
            self._process_song_file(song_file, image_files[0] if image_files else None)
        else:
            # Process images only (no text files allowed in non-song directories)
            for image_file in image_files:
                self.add_slide_with_background(str(image_file))
        
        return self.prs
    
    def _process_song_file(self, song_file: Path, background_image: Optional[Path]):
        """Process a song file with arrangement."""
        # Parse song file (similar to PP6 generator)
        sections, arrangement = self._parse_song_file(str(song_file))
        
        # Process each section in arrangement order
        for section_name in arrangement:
            if section_name not in sections:
                continue
            
            section_lines = sections[section_name]
            
            # Split section into pages
            for i in range(0, len(section_lines), self.page_break_every):
                page_lines = section_lines[i:i + self.page_break_every]
                text = '\n'.join(line.strip() for line in page_lines if line.strip())
                
                # Create slide
                self.add_slide_with_background(
                    str(background_image) if background_image else None,
                    text
                )
    
    def _parse_song_file(self, filepath: str) -> Tuple[Dict[str, List[str]], List[str]]:
        """Parse song file to extract sections and arrangement."""
        with open(filepath, 'r', encoding='utf-8') as f:
            lines = f.read().strip().split('\n')
        
        # Find arrangement
        arrangement = []
        arrangement_line_index = -1
        for i, line in enumerate(lines):
            if line.strip().startswith('Arrangement'):
                arrangement_line_index = i
                if i + 1 < len(lines):
                    arrangement = lines[i + 1].strip().split()
                break
        
        # Extract sections
        valid_sections = set(arrangement)
        sections = {}
        current_section = None
        
        for i, line in enumerate(lines):
            if i == arrangement_line_index or i == arrangement_line_index + 1:
                continue
                
            line_stripped = line.strip()
            
            if line_stripped in valid_sections:
                current_section = line_stripped
                sections[current_section] = []
            elif current_section:
                sections[current_section].append(line)
        
        return sections, arrangement
    
    def save(self, output_path: str):
        """Save the presentation to file."""
        self.prs.save(output_path)
        print(f"Generated PowerPoint presentation: {output_path}")
        print(f"  Total slides: {len(self.prs.slides)}")